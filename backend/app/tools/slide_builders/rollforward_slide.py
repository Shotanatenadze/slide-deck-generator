"""Portfolio Roll Forward slide — two-column table (Description | Amount)."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.tools.slide_builders import (
    ACCENT_RED, DARK_GRAY, FONT_NAME, LIGHT_GRAY, NAVY, PRIMARY_BLUE, WHITE,
)


def build(prs, data: dict, commentary: str | None = None) -> None:
    """
    Add a Portfolio Roll Forward slide.

    data keys:
      - roll_forward: list[dict] with keys label, value
    """
    roll_forward = data.get("roll_forward", [])
    if not roll_forward:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "PORTFOLIO ROLL FORWARD"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # Table
    rows = len(roll_forward) + 1  # header row + data
    cols = 2
    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(2.0), Inches(1.5),
        Inches(9.0), Inches(0.40 * rows),
    )
    table = tbl_shape.table

    table.columns[0].width = Inches(5.5)
    table.columns[1].width = Inches(3.5)

    # Header
    for ci, h in enumerate(["Description", "Amount"]):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, rf in enumerate(roll_forward):
        row_idx = ri + 1
        label = rf.get("label", "")
        value = rf.get("value")
        is_total = "ending" in label.lower() or "total" in label.lower()
        is_negative = value is not None and value < 0

        # Format value
        if value is not None:
            val_str = f"${value:,.2f}" if value >= 0 else f"(${abs(value):,.2f})"
        else:
            val_str = "---"

        for ci, text in enumerate([label, val_str]):
            cell = table.cell(row_idx, ci)
            cell.text = text

            # Styling
            if is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_BLUE
                font_color = WHITE
                bold = True
            elif ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
                font_color = ACCENT_RED if (ci == 1 and is_negative) else DARK_GRAY
                bold = False
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                font_color = ACCENT_RED if (ci == 1 and is_negative) else DARK_GRAY
                bold = False

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.bold = bold
                paragraph.font.color.rgb = font_color
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.8), Inches(12), Inches(0.3),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA PORTFOLIO ROLLFORWARD REPORT IN CLEARWATER"
    fp.font.size = Pt(7)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True
