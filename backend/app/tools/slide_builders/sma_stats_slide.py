"""SMA Portfolio Statistics slide — summary stats table + sector & credit pie charts."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.tools.slide_builders import (
    DARK_GRAY, FONT_NAME, LIGHT_GRAY, NAVY, PRIMARY_BLUE, WHITE,
)
from app.tools.chart_builder import add_pie_chart


def build(prs, data: dict, commentary: str | None = None) -> None:
    """
    Add an SMA Portfolio Statistics slide.

    data keys:
      - sma_summary: dict with keys total_market_value, avg_duration, avg_yield,
        num_holdings, sector_allocation (dict), credit_quality (dict)
    """
    summary = data.get("sma_summary")
    if not summary:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "SMA PORTFOLIO STATISTICS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # --- Summary stats table ---
    stats = [
        ("Total Market Value", f"${(summary.get('total_market_value') or 0):,.0f}"),
        ("Number of Holdings", str(summary.get("num_holdings") or 0)),
        ("Average Duration", f"{(summary.get('avg_duration') or 0):.2f} yrs"),
        ("Average Yield", f"{(summary.get('avg_yield') or 0):.2f}%"),
    ]
    if summary.get("weighted_sp_rating"):
        stats.append(("Weighted S&P Rating", summary["weighted_sp_rating"]))
    if summary.get("weighted_moody_rating"):
        stats.append(("Weighted Moody's Rating", summary["weighted_moody_rating"]))

    rows = len(stats) + 1
    tbl_shape = slide.shapes.add_table(
        rows, 2,
        Inches(0.4), Inches(1.2),
        Inches(4.5), Inches(0.35 * rows),
    )
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)

    # Header
    for ci, h in enumerate(["Statistic", "Value"]):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, (label, val) in enumerate(stats):
        row_idx = ri + 1
        for ci, text in enumerate([label, val]):
            cell = table.cell(row_idx, ci)
            cell.text = text
            bg = LIGHT_GRAY if ri % 2 == 1 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Sector Allocation pie chart ---
    sector_alloc = summary.get("sector_allocation", {})
    if sector_alloc:
        # Convert percentage values to fractions for the chart
        pie_data = {k: v / 100.0 for k, v in sector_alloc.items() if v > 0}
        if pie_data:
            add_pie_chart(
                slide, pie_data,
                left=5.2, top=1.0, width=3.8, height=3.5,
                title="Sector Allocation",
            )

    # --- Credit Quality pie chart ---
    credit_q = summary.get("credit_quality", {})
    if credit_q:
        pie_data = {k: v / 100.0 for k, v in credit_q.items() if v > 0}
        if pie_data:
            add_pie_chart(
                slide, pie_data,
                left=9.2, top=1.0, width=3.8, height=3.5,
                title="Credit Quality",
            )

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.8), Inches(12), Inches(0.3),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA SMA HOLDINGS REPORT IN CLEARWATER ANALYTICS"
    fp.font.size = Pt(7)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True
