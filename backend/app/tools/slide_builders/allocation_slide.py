"""Asset Allocation slide — table on the left, pie chart on the right."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from app.tools.slide_builders import (
    DARK_GRAY, FONT_NAME, LIGHT_GRAY, NAVY, PRIMARY_BLUE, WHITE,
)
from app.tools.chart_builder import add_pie_chart


def build(prs, data: dict, commentary: str | None = None) -> None:
    """
    Add an Asset Allocation slide.

    data keys:
      - allocation: list[dict] with keys strategy, market_value, actual_pct, target_pct
    """
    allocation = data.get("allocation", [])
    if not allocation:
        return

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # --- Header bar ---
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "ASSET ALLOCATION"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # --- Table ---
    rows = len(allocation) + 2  # header + data rows + total row
    cols = 4
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.4), Inches(1.2),
        Inches(7.2), Inches(0.35 * rows),
    )
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.3)
    table.columns[3].width = Inches(1.3)

    # Header row
    headers = ["STRATEGY", "MARKET VALUE", "ACTUAL %", "TARGET %"]
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    total_mv = 0.0
    for ri, alloc in enumerate(allocation):
        row_idx = ri + 1
        mv = alloc.get("market_value") or 0
        total_mv += mv
        actual = alloc.get("actual_pct")
        target = alloc.get("target_pct")

        values = [
            alloc.get("strategy", ""),
            f"${mv:,.0f}",
            f"{actual * 100:.1f}%" if actual is not None else "---",
            f"{target * 100:.1f}%" if target is not None else "---",
        ]

        for ci, val in enumerate(values):
            cell = table.cell(row_idx, ci)
            cell.text = val
            # Alternating row shading
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Total row
    total_row = rows - 1
    total_values = ["TOTAL", f"${total_mv:,.0f}", "100.0%", ""]
    for ci, val in enumerate(total_values):
        cell = table.cell(total_row, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Pie chart ---
    pie_data = {}
    for alloc in allocation:
        strategy = alloc.get("strategy", "")
        pct = alloc.get("actual_pct")
        if pct is not None and pct > 0:
            pie_data[strategy] = pct
    if not pie_data:
        # Fallback: use market value proportions
        for alloc in allocation:
            strategy = alloc.get("strategy", "")
            mv = alloc.get("market_value") or 0
            if mv > 0 and total_mv > 0:
                pie_data[strategy] = mv / total_mv

    if pie_data:
        add_pie_chart(
            slide,
            pie_data,
            left=8.0, top=1.2, width=4.8, height=4.5,
            title="Allocation Breakdown",
        )

    # Commentary (if provided)
    if commentary:
        ctb = slide.shapes.add_textbox(
            Inches(0.4), Inches(6.0), Inches(7.2), Inches(0.6),
        )
        ctf = ctb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = commentary
        cp.font.size = Pt(9)
        cp.font.italic = True
        cp.font.color.rgb = DARK_GRAY
        cp.font.name = FONT_NAME

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.8), Inches(12), Inches(0.3),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA BOARD REPORT IN CLEARWATER ANALYTICS"
    fp.font.size = Pt(7)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True
