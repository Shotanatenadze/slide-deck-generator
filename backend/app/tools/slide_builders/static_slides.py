"""Static slides: section dividers and disclaimer."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.tools.slide_builders import DARK_GRAY, FONT_NAME, NAVY, PRIMARY_BLUE, WHITE


def build_section_divider(prs, title: str, subtitle: str = "") -> None:
    """Add a branded section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full-slide background
    bg = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Accent bar on left
    bar = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(0.15), prs.slide_height,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()

    # Title
    txbox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xBB, 0xCE, 0xDA)
        p2.font.name = FONT_NAME
        p2.space_before = Pt(12)


def build_disclaimer(prs) -> None:
    """Add a disclaimer slide at the end of the deck."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "DISCLAIMERS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    disclaimer_text = (
        "This report is generated using a third-party proprietary system and is provided "
        "to you as an accommodation in the review of your investment activity. Although the "
        "informational sources used are deemed to be reliable, the firm does not guarantee or make "
        "any warranty regarding the accuracy, suitability or completeness of such information.\n\n"
        "The data should not be used in lieu of statements received directly from your custodian. "
        "The firm does not provide tax advice, or legal advice through this site or otherwise.\n\n"
        "The contents in this report are for informational purposes only and intended solely for "
        "the user and should not be distributed to any third party without prior written consent. "
        "Opinions, estimates and assumptions expressed herein reflect our judgment as of the date "
        "of publication and are subject to change without notice.\n\n"
        "Any statements regarding performance may not be realized and past performance is not "
        "indicative of future results. Investors should note that the value of any investment "
        "strategy or security may fluctuate and underlying principal values may rise or fall.\n\n"
        "[Insert applicable regulatory disclosures here.]"
    )

    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = disclaimer_text
    p.font.size = Pt(9)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_NAME
    p.line_spacing = Pt(14)
