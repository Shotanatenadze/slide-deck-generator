"""Title slide — company name, client name, quarter/year, branded header bar."""

from __future__ import annotations

from datetime import datetime

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.tools.slide_builders import DARK_GRAY, FONT_NAME, NAVY, PRIMARY_BLUE, WHITE


def build(prs, data: dict) -> None:
    """
    Add a title slide.

    data keys:
      - client_name: str
      - as_of_date: str | None
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        prs.slide_width, Inches(1.4),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    # Company name
    txbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(10), Inches(0.6))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = data.get("firm_name", "INVESTMENT ADVISORY")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(10), Inches(0.4))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = data.get("firm_subtitle", "Portfolio Management")
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xBB, 0xCE, 0xDA)
    p2.font.name = FONT_NAME

    # Client name
    client_name = data.get("client_name", "Client")
    ctb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(1.0))
    tf3 = ctb.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = client_name
    p3.font.size = Pt(32)
    p3.font.bold = True
    p3.font.color.rgb = NAVY
    p3.font.name = FONT_NAME

    # Quarter / date line
    as_of = data.get("as_of_date") or ""
    deck_title = data.get("deck_title") or ""
    if deck_title:
        if as_of:
            date_text = f"{deck_title}  |  As of {as_of}"
        else:
            now = datetime.utcnow()
            quarter = (now.month - 1) // 3 + 1
            date_text = f"{deck_title}  |  Q{quarter} {now.year}"
    elif as_of:
        date_text = f"Quarterly Investment Review  |  As of {as_of}"
    else:
        now = datetime.utcnow()
        quarter = (now.month - 1) // 3 + 1
        date_text = f"Quarterly Investment Review  |  Q{quarter} {now.year}"

    dtb = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11), Inches(0.5))
    tf4 = dtb.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = date_text
    p4.font.size = Pt(16)
    p4.font.color.rgb = PRIMARY_BLUE
    p4.font.name = FONT_NAME

    # Thin accent line
    line = slide.shapes.add_shape(
        1,
        Inches(0.8), Inches(4.1),
        Inches(4), Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_BLUE
    line.line.fill.background()

    # Footer
    ftb = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11), Inches(0.4))
    tf5 = ftb.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "CONFIDENTIAL — FOR CLIENT USE ONLY"
    p5.font.size = Pt(9)
    p5.font.color.rgb = DARK_GRAY
    p5.font.name = FONT_NAME
    p5.alignment = PP_ALIGN.LEFT
