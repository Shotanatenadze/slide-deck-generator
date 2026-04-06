"""Performance Summary slide — clustered bar chart comparing Net Return vs Index."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.tools.slide_builders import DARK_GRAY, FONT_NAME, NAVY, WHITE
from app.tools.chart_builder import add_bar_chart

# Short labels for performance periods
PERIOD_SHORT = {
    "Quarter to Date": "QTD",
    "Year to Date": "YTD",
    "Trailing Year": "1 Year",
    "Trailing 3 Years": "3 Year",
    "Trailing 5 Years": "5 Year",
    "Trailing 10 Years": "10 Year",
    "Since Inception": "Inception",
}


def build(prs, data: dict, commentary: str | None = None) -> None:
    """
    Add a Performance Summary slide.

    data keys:
      - performance: list[dict] with keys period, total_return, index_return
    """
    performance = data.get("performance", [])
    if not performance:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "PERFORMANCE SUMMARY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # Build chart data — skip periods where both return values are None
    categories: list[str] = []
    net_values: list[float] = []
    index_values: list[float] = []

    for row in performance:
        period = row.get("period", "")
        net = row.get("total_return")
        idx = row.get("index_return")
        if net is None and idx is None:
            continue
        label = PERIOD_SHORT.get(period, period)
        categories.append(label)
        net_values.append(net if net is not None else 0)
        index_values.append(idx if idx is not None else 0)

    if not categories:
        return

    series_list = [
        {"name": "Net Return", "values": net_values},
        {"name": "Index Return", "values": index_values},
    ]

    add_bar_chart(
        slide,
        categories,
        series_list,
        left=0.8, top=1.2, width=11.5, height=5.0,
        title="Net Performance (Includes Oversight Fee)",
    )

    # Performance table below/beside chart — small summary table
    rows = len(categories) + 1
    cols = 3
    tbl_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.2),
        Inches(5.0), Inches(0.28 * rows),
    )
    # Make table semi-transparent by hiding it behind chart area
    # Actually, place it to the right or below — let's put a compact version below
    # We'll skip the table when chart is present for cleanliness.
    # Remove the table shape — chart alone is the primary visual
    sp = tbl_shape._element
    sp.getparent().remove(sp)

    # Commentary
    if commentary:
        ctb = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
        )
        ctf = ctb.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = commentary
        cp.font.size = Pt(9)
        cp.font.italic = True
        cp.font.color.rgb = DARK_GRAY
        cp.font.name = FONT_NAME

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.8), Inches(12), Inches(0.3),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA BOARD REPORT IN CLEARWATER ANALYTICS"
    fp.font.size = Pt(7)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True
