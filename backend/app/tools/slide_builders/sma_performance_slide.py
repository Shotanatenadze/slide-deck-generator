"""SMA Performance slide — performance bar chart for SMA strategies."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.tools.slide_builders import DARK_GRAY, FONT_NAME, NAVY, WHITE
from app.tools.chart_builder import add_bar_chart

PERIOD_SHORT = {
    "Quarter to Date": "QTD",
    "Year to Date": "YTD",
    "Trailing Year": "1 Year",
    "Trailing 3 Years": "3 Year",
    "Trailing 5 Years": "5 Year",
    "Trailing 10 Years": "10 Year",
    "Since Inception": "Inception",
}


def build(prs, data: dict) -> None:
    """
    Add an SMA Performance slide.

    data keys:
      - performance: list[dict] with keys period, total_return, index_return
        (SMA-specific performance data)
    """
    performance = data.get("performance", [])
    if not performance:
        return

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.8),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
    p = htb.text_frame.paragraphs[0]
    p.text = "SMA PERFORMANCE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # Build chart data
    categories: list[str] = []
    net_values: list[float] = []
    index_values: list[float] = []

    for row in performance:
        period = row.get("period", "")
        net = row.get("total_return")
        idx = row.get("index_return")
        if net is None and idx is None:
            continue
        label = PERIOD_SHORT.get(period, period)
        categories.append(label)
        net_values.append(net if net is not None else 0)
        index_values.append(idx if idx is not None else 0)

    if not categories:
        return

    series_list = [
        {"name": "Net Return", "values": net_values},
        {"name": "Index Return", "values": index_values},
    ]

    add_bar_chart(
        slide,
        categories,
        series_list,
        left=0.8, top=1.2, width=11.5, height=5.0,
        title="SMA Fixed Income Performance (Net of Fees)",
    )

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.8), Inches(12), Inches(0.3),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA BOARD REPORT IN CLEARWATER ANALYTICS"
    fp.font.size = Pt(7)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True
