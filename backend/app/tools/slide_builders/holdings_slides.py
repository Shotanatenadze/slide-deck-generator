"""SMA Holdings slides — paginated bond holdings tables (25 rows per page)."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.tools.slide_builders import (
    DARK_GRAY, FONT_NAME, LIGHT_GRAY, NAVY, PRIMARY_BLUE, WHITE,
)

ROWS_PER_PAGE = 25

# Column definitions: (header, key, width_inches, alignment)
COLUMNS = [
    ("ISSUER", "issuer", 2.0, PP_ALIGN.LEFT),
    ("SECTOR", "sector", 1.5, PP_ALIGN.LEFT),
    ("DURATION", "duration", 1.0, PP_ALIGN.RIGHT),
    ("YIELD", "yield_to_worst", 0.8, PP_ALIGN.RIGHT),
    ("S&P", "sp_rating", 0.7, PP_ALIGN.CENTER),
    ("MOODY'S", "moody_rating", 0.8, PP_ALIGN.CENTER),
    ("MATURITY", "maturity_date", 1.2, PP_ALIGN.CENTER),
    ("MKT VALUE", "market_value", 1.5, PP_ALIGN.RIGHT),
    ("% PORT", "pct_of_portfolio", 0.8, PP_ALIGN.RIGHT),
]


def build(prs, data: dict) -> None:
    """
    Add paginated SMA Holdings slides.

    data keys:
      - sma_holdings: list[dict] with SMAHolding fields
    """
    holdings = data.get("sma_holdings") or []
    if not holdings:
        return

    # Paginate
    pages = [
        holdings[i : i + ROWS_PER_PAGE]
        for i in range(0, len(holdings), ROWS_PER_PAGE)
    ]

    for page_num, page_holdings in enumerate(pages, start=1):
        _build_page(prs, page_holdings, page_num, len(pages))


def _build_page(prs, holdings: list[dict], page_num: int, total_pages: int) -> None:
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0),
        prs.slide_width, Inches(0.7),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()

    title_text = "SMA HOLDINGS"
    if total_pages > 1:
        title_text += f"  ({page_num}/{total_pages})"

    htb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(10), Inches(0.45))
    p = htb.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME

    # Table
    num_cols = len(COLUMNS)
    num_rows = len(holdings) + 1  # header + data
    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.2), Inches(0.9),
        Inches(12.9), Inches(0.22 * num_rows),
    )
    table = tbl_shape.table

    # Column widths
    for ci, (_, _, w, _) in enumerate(COLUMNS):
        table.columns[ci].width = Inches(w)

    # Header row
    for ci, (hdr, _, _, align) in enumerate(COLUMNS):
        cell = table.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(7)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = align
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for ri, holding in enumerate(holdings):
        row_idx = ri + 1
        for ci, (_, key, _, align) in enumerate(COLUMNS):
            cell = table.cell(row_idx, ci)
            raw = holding.get(key)
            text = _format_cell(key, raw)
            cell.text = text

            # Alternating row shading
            bg = LIGHT_GRAY if ri % 2 == 1 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(7)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.font.name = FONT_NAME
                paragraph.alignment = align
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Source footnote
    ftb = slide.shapes.add_textbox(
        Inches(0.3), Inches(6.9), Inches(12), Inches(0.25),
    )
    fp = ftb.text_frame.paragraphs[0]
    fp.text = "SOURCED VIA SMA HOLDINGS REPORT IN CLEARWATER ANALYTICS"
    fp.font.size = Pt(6)
    fp.font.color.rgb = DARK_GRAY
    fp.font.name = FONT_NAME
    fp.font.italic = True


def _format_cell(key: str, value) -> str:
    """Format a cell value based on the column type."""
    if value is None:
        return "---"
    if key == "duration":
        return f"{value:.2f}" if isinstance(value, (int, float)) else str(value)
    if key == "yield_to_worst":
        return f"{value:.2f}%" if isinstance(value, (int, float)) else str(value)
    if key == "market_value":
        return f"${value:,.0f}" if isinstance(value, (int, float)) else str(value)
    if key == "pct_of_portfolio":
        if isinstance(value, (int, float)):
            return f"{value * 100:.2f}%" if value < 1 else f"{value:.2f}%"
        return str(value)
    if key == "maturity_date":
        # Already ISO string or None
        return str(value)[:10] if value else "---"
    return str(value) if value else "---"
