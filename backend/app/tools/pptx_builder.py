"""
High-level deck assembly.

Orchestrates slide builders to produce a complete .pptx file.
"""

from __future__ import annotations

import os
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu

from app.config import settings
from app.tools.slide_builders import title_slide
from app.tools.slide_builders import allocation_slide
from app.tools.slide_builders import performance_slide
from app.tools.slide_builders import rollforward_slide
from app.tools.slide_builders import sma_stats_slide
from app.tools.slide_builders import sma_performance_slide
from app.tools.slide_builders import holdings_slides
from app.tools.slide_builders import static_slides

# Widescreen 13.333" x 7.5"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def build_deck(
    portfolio_data: dict,
    market_data: dict | None = None,
    compliance_report: dict | None = None,
    analyst_prompt: str = "",
    generation_id: str | None = None,
    commentary: dict | None = None,
) -> str:
    """
    Build a complete PowerPoint deck from structured data.

    Parameters
    ----------
    portfolio_data : dict
        Parsed portfolio data (PortfolioData-compatible dict).
    market_data : dict | None
        Market context narrative sections.
    compliance_report : dict | None
        Compliance check results (included as metadata, not as a slide).
    analyst_prompt : str
        Free-text analyst instructions (used for future customization).
    generation_id : str | None
        Used for file naming.

    Returns
    -------
    str : Absolute path to the generated .pptx file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    client_name = portfolio_data.get("client_name", "Client")

    c = commentary or {}

    # ----------------------------------------------------------------
    # 1. Title slide
    # ----------------------------------------------------------------
    title_slide.build(prs, {
        "client_name": client_name,
        "as_of_date": portfolio_data.get("as_of_date"),
        "deck_title": c.get("deck_title"),
    })

    # ----------------------------------------------------------------
    # 2. Agenda / section divider
    # ----------------------------------------------------------------
    static_slides.build_section_divider(
        prs, "Investment Report",
        subtitle="Portfolio Review & Performance Analysis",
    )

    # ----------------------------------------------------------------
    # 3. Market Update section (if market data provided)
    # ----------------------------------------------------------------
    if market_data and market_data.get("sections"):
        static_slides.build_section_divider(prs, "Market Update")
        _build_market_slides(prs, market_data)

    # ----------------------------------------------------------------
    # 4. Asset Allocation (table + pie chart)
    # ----------------------------------------------------------------
    if portfolio_data.get("allocation"):
        allocation_slide.build(prs, portfolio_data, c.get("allocation_commentary"))

    # ----------------------------------------------------------------
    # 5. Performance (bar chart)
    # ----------------------------------------------------------------
    if portfolio_data.get("performance"):
        performance_slide.build(prs, portfolio_data, c.get("performance_commentary"))

    # ----------------------------------------------------------------
    # 6. Roll Forward table
    # ----------------------------------------------------------------
    if portfolio_data.get("roll_forward"):
        rollforward_slide.build(prs, portfolio_data, c.get("rollforward_commentary"))

    # ----------------------------------------------------------------
    # 7. SMA sections (if SMA data is present)
    # ----------------------------------------------------------------
    if portfolio_data.get("sma_summary") or portfolio_data.get("sma_holdings"):
        static_slides.build_section_divider(
            prs, "Fixed Income SMA",
            subtitle="Separately Managed Account Detail",
        )

        # SMA Statistics + pie charts
        if portfolio_data.get("sma_summary"):
            sma_stats_slide.build(prs, portfolio_data, c.get("sma_commentary"))

        # SMA Performance
        if portfolio_data.get("performance"):
            sma_performance_slide.build(prs, portfolio_data)

        # SMA Holdings (paginated)
        if portfolio_data.get("sma_holdings"):
            holdings_slides.build(prs, portfolio_data)

    # ----------------------------------------------------------------
    # 8. Disclaimer
    # ----------------------------------------------------------------
    static_slides.build_disclaimer(prs)

    # ----------------------------------------------------------------
    # Save
    # ----------------------------------------------------------------
    gen_id = generation_id or uuid.uuid4().hex[:12]
    out_dir = Path(settings.GENERATED_DIR) / gen_id
    out_dir.mkdir(parents=True, exist_ok=True)

    filename = f"{client_name.replace(' ', '_')}_deck.pptx"
    out_path = out_dir / filename
    prs.save(str(out_path))

    return str(out_path.resolve())


def _clean_market_text(text: str) -> str:
    """Strip markdown links, URLs, and table formatting from market text."""
    import re
    # Convert markdown links [text](url) to just text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    # Remove bare URLs
    text = re.sub(r'https?://\S+', '', text)
    # Remove markdown bold/italic
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # Remove markdown table separators
    text = re.sub(r'\|[\s-]+\|[\s-]+\|[\s-]+\|', '', text)
    # Clean up markdown table rows into readable text
    text = re.sub(r'^\s*\|(.+)\|\s*$', lambda m: m.group(1).replace('|', ' — ').strip(), text, flags=re.MULTILINE)
    # Collapse multiple blank lines
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _build_market_slides(prs, market_data: dict) -> None:
    """Build professional market update slides — one per section."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from app.tools.slide_builders import DARK_GRAY, FONT_NAME, NAVY, PRIMARY_BLUE, WHITE

    sections = market_data.get("sections", [])
    if not sections:
        # Single slide from raw text
        raw = _clean_market_text(market_data.get("raw_text", ""))
        if raw:
            sections = [{"title": "Market Overview", "content": raw}]

    for section in sections:
        title = section.get("title", "Market Update")
        content = _clean_market_text(section.get("content", ""))
        if not content:
            continue

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Header bar
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            prs.slide_width, Inches(0.8),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = NAVY
        header.line.fill.background()

        htb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.5))
        p = htb.text_frame.paragraphs[0]
        p.text = f"MARKET UPDATE — {title.upper()}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_NAME

        # Content body — split into paragraphs for clean formatting
        txbox = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5),
        )
        tf = txbox.text_frame
        tf.word_wrap = True

        paragraphs = [p.strip() for p in content.split("\n") if p.strip()]
        for i, para_text in enumerate(paragraphs[:15]):  # Cap at 15 paragraphs per slide
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = para_text[:500]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.font.name = FONT_NAME
            p.space_after = Pt(6)
