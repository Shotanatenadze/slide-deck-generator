"""
Native PPTX chart utilities.

Creates pie and bar charts using python-pptx's chart API.
All charts use branding colors and Poppins / Calibri fonts.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Branding colors
# ---------------------------------------------------------------------------
PRIMARY_BLUE = RGBColor(0x10, 0x62, 0x80)
DARK_GRAY = RGBColor(0x39, 0x48, 0x49)
NAVY = RGBColor(0x0F, 0x2E, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0x4E, 0xA8, 0xC8)
ACCENT_GREEN = RGBColor(0x5B, 0xB5, 0x73)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x39)
ACCENT_RED = RGBColor(0xD9, 0x53, 0x4F)
ACCENT_PURPLE = RGBColor(0x7B, 0x68, 0xAF)
ACCENT_TEAL = RGBColor(0x2E, 0x9E, 0x9E)

PIE_COLORS = [
    PRIMARY_BLUE, LIGHT_BLUE, ACCENT_GREEN, ACCENT_ORANGE,
    ACCENT_RED, ACCENT_PURPLE, ACCENT_TEAL, DARK_GRAY,
    NAVY, RGBColor(0xA0, 0xA0, 0xA0),
]

BAR_SERIES_COLORS = [PRIMARY_BLUE, ACCENT_GREEN, ACCENT_ORANGE, DARK_GRAY]

FONT_NAME = "Calibri"  # Poppins preferred, Calibri as safe fallback


def add_pie_chart(
    slide,
    data_dict: dict[str, float],
    left: float,
    top: float,
    width: float,
    height: float,
    colors: list[RGBColor] | None = None,
    title: str | None = None,
) -> None:
    """
    Add a pie chart to a slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
    data_dict : {label: value, ...}
    left, top, width, height : inches
    colors : optional list of RGBColor per slice
    title : optional chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = list(data_dict.keys())
    chart_data.add_series("Allocation", list(data_dict.values()))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_frame.chart

    # Title
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.name = FONT_NAME

    # Data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0.0%"
    plot.data_labels.font.size = Pt(8)
    plot.data_labels.font.name = FONT_NAME
    plot.data_labels.font.color.rgb = DARK_GRAY

    # Slice colors
    use_colors = colors or PIE_COLORS
    series = plot.series[0]
    for idx in range(len(data_dict)):
        point = series.points[idx]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = use_colors[idx % len(use_colors)]


def add_bar_chart(
    slide,
    categories: list[str],
    series_list: list[dict],
    left: float,
    top: float,
    width: float,
    height: float,
    title: str | None = None,
) -> None:
    """
    Add a clustered bar chart to a slide.

    Parameters
    ----------
    slide : pptx.slide.Slide
    categories : list of category labels (e.g. period names)
    series_list : list of {"name": str, "values": list[float|None]}
        Each dict is one bar series.
    left, top, width, height : inches
    title : optional chart title
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series_list:
        # Replace None with 0 for chart data
        values = [v if v is not None else 0 for v in s["values"]]
        chart_data.add_series(s["name"], values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_frame.chart

    # Title
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.chart_title.text_frame.paragraphs[0].font.name = FONT_NAME
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = FONT_NAME

    # Value axis formatting
    value_axis = chart.value_axis
    value_axis.has_title = False
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    value_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.name = FONT_NAME
    value_axis.tick_labels.number_format = "0.0%"

    # Category axis formatting
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(8)
    cat_axis.tick_labels.font.name = FONT_NAME
    cat_axis.format.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)

    # Series colors
    plot = chart.plots[0]
    plot.gap_width = 100
    for idx, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = BAR_SERIES_COLORS[idx % len(BAR_SERIES_COLORS)]
